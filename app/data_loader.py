from langchain_community.document_loaders.powerpoint import UnstructuredPowerPointLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
from PIL import Image
import pytesseract
import os
import io


pytesseract.pytesseract.tesseract_cmd = r"C:Program Files/Tesseract-OCR/tesseract.exe"

def load_ppt_text(filepath: str):
    loader = UnstructuredPowerPointLoader(filepath)
    return loader.load()

def extract_images_text(filepath: str, output_dir="extracted_images"):
    prs = Presentation(filepath)
    os.makedirs(output_dir, exist_ok=True)

    ocr_texts = []
    image_count = 0

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  
                image = shape.image
                image_bytes = image.blob
                ext = image.ext

                image_path = os.path.join(output_dir, f"slide{slide_num}_img{image_count}.{ext}")
                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                pil_image = Image.open(io.BytesIO(image_bytes))
                text = pytesseract.image_to_string(pil_image)
                if text.strip():
                    ocr_texts.append(f"[Slide {slide_num} - Image {image_count}] {text}")

                image_count += 1

    return ocr_texts

def load_ppt(filepath: str):
    text_docs = load_ppt_text(filepath)
    ocr_texts = extract_images_text(filepath)

    all_texts = [doc.page_content for doc in text_docs] + ocr_texts
    return all_texts

def split_text(texts):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=20)
    return splitter.create_documents(texts)

